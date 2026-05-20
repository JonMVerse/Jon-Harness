"""
Multiverse branded slide builder — helper library for python-pptx.

Usage: import or copy into your generation script, then call the helper
functions. All brand colours and layout constants are defined here so
generation scripts stay clean.

Brand source of truth: Drive file 19yCvhpI6TnCD1J_Afm4zNlJzTvDLjdZKERf62JV5fOg
Hex values here are calibrated approximations — verify against the brand deck
if exact fidelity is required.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Brand colours ────────────────────────────────────────────────────────────

ULTRAVIOLET = RGBColor(0x52, 0x00, 0xF2)   # Primary brand purple — covers/dividers
LUNAR_50    = RGBColor(0xF5, 0xF4, 0xF0)   # Light cream — PRIMARY background
LUNAR_700   = RGBColor(0x70, 0x6E, 0x6A)   # Medium grey — eyebrows, body text
BLACKHOLE   = RGBColor(0x14, 0x14, 0x0F)   # Near-black — titles on Lunar
SUNBEAM     = RGBColor(0xFF, 0xD8, 0x15)   # Yellow — accent bar, title highlight
STARLIGHT   = RGBColor(0xF0, 0xEE, 0xE9)  # Near-white — alternative light bg
ECLIPSE     = RGBColor(0x1A, 0x00, 0x57)  # Dark purple — leadership/closing slides
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# Tinted purple for subtitle text on dark (Ultraviolet / Eclipse) slides
PURPLE_TINT = RGBColor(0xCC, 0xBB, 0xFF)
PURPLE_DIM  = RGBColor(0xAA, 0x99, 0xDD)

# Secondary palette — data/charts only
DARK_ULTRAVIOLET = RGBColor(0x3A, 0x00, 0xB4)
DARK_SUNBEAM     = RGBColor(0xD4, 0xB0, 0x00)
GALAXY           = RGBColor(0x7B, 0x5E, 0xF8)
SOFT_GALAXY      = RGBColor(0xAA, 0x96, 0xFA)
LUNAR_600        = RGBColor(0x8A, 0x88, 0x84)
DARK_ECLIPSE     = RGBColor(0x0D, 0x00, 0x33)


# ── Layout constants ─────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)   # Google Slides widescreen
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.6)
MARGIN_R = Inches(12.73)  # MARGIN_L + content width
CONTENT_W = Inches(12.13)
CONTENT_T = Inches(1.55)
CONTENT_H = Inches(5.5)


# ── Core helpers ─────────────────────────────────────────────────────────────

def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, colour: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _box(slide, text: str, l, t, w, h,
         size=12, bold=False, colour=WHITE,
         align=PP_ALIGN.LEFT, wrap=True,
         font="Inter") -> object:
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name  = font
    f.size  = Pt(size)
    f.bold  = bold
    f.color.rgb = colour
    return tb


def _accent_bar(slide, top=Inches(0.18)):
    s = slide.shapes.add_shape(1, MARGIN_L, top, Inches(0.55), Inches(0.06))
    s.fill.solid()
    s.fill.fore_color.rgb = SUNBEAM
    s.line.fill.background()


def _eyebrow(slide, text: str, top=Inches(0.40)):
    _box(slide, text.upper(), MARGIN_L, top, Inches(8), Inches(0.35),
         size=9, colour=LUNAR_700)


def _divider_line(slide, x, t, h):
    s = slide.shapes.add_shape(1, x, t, Inches(0.02), h)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0xD8, 0xD6, 0xD2)
    s.line.fill.background()


# ── Slide-type factories ──────────────────────────────────────────────────────

def make_cover(prs, title: str, subtitle: str = "",
               footer: str = "", teams: str = ""):
    """Ultraviolet cover slide."""
    slide = _blank(prs)
    _bg(slide, ULTRAVIOLET)

    # Sunbeam bar
    bar = slide.shapes.add_shape(1, MARGIN_L, Inches(0.55), Inches(0.55), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = SUNBEAM; bar.line.fill.background()

    _box(slide, title,
         MARGIN_L, Inches(1.6), Inches(11), Inches(1.8),
         size=46, bold=True, colour=WHITE)

    if subtitle:
        _box(slide, subtitle,
             MARGIN_L, Inches(3.2), Inches(8), Inches(0.7),
             size=24, colour=WHITE)

    if teams:
        _box(slide, teams,
             MARGIN_L, Inches(3.9), Inches(11), Inches(0.6),
             size=14, colour=PURPLE_TINT)

    if footer:
        _box(slide, footer,
             MARGIN_L, Inches(6.6), Inches(6), Inches(0.5),
             size=10, colour=PURPLE_TINT)


def make_divider(prs, team_name: str, section: str = "",
                 note: str = "", bg=None):
    """
    Ultraviolet section divider with large team/section name.
    bg defaults to ULTRAVIOLET; pass Eclipse for leadership-style dividers.
    """
    bg = bg or ULTRAVIOLET
    slide = _blank(prs)
    _bg(slide, bg)

    if section:
        _box(slide, section,
             MARGIN_L, Inches(0.45), Inches(2), Inches(0.4),
             size=11, colour=PURPLE_TINT)

    _box(slide, team_name,
         MARGIN_L, Inches(2.5), Inches(11), Inches(1.5),
         size=56, bold=True, colour=WHITE)

    if note:
        _box(slide, note,
             MARGIN_L, Inches(4.1), Inches(8), Inches(0.5),
             size=13, colour=PURPLE_TINT)


def make_leadership_divider(prs, title: str = "Leadership",
                             subtitle: str = "Direction  ·  Kudos  ·  One takeaway",
                             note: str = "", presenter: str = "Yuval"):
    """Eclipse-background leadership/closing divider."""
    slide = _blank(prs)
    _bg(slide, ECLIPSE)

    bar = slide.shapes.add_shape(1, MARGIN_L, Inches(0.55), Inches(0.55), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = SUNBEAM; bar.line.fill.background()

    _box(slide, title,
         MARGIN_L, Inches(2.0), Inches(11), Inches(1.0),
         size=46, bold=True, colour=WHITE)

    if subtitle:
        _box(slide, subtitle,
             MARGIN_L, Inches(3.1), Inches(10), Inches(0.55),
             size=14, colour=PURPLE_TINT)

    if note:
        _box(slide, note,
             MARGIN_L, Inches(3.75), Inches(10), Inches(0.45),
             size=11, colour=PURPLE_DIM)

    if presenter:
        _box(slide, presenter,
             MARGIN_L, Inches(6.6), Inches(4), Inches(0.45),
             size=10, colour=PURPLE_TINT)


def make_content(prs, eyebrow: str = "", title: str = "",
                 body: str = "", note: str = ""):
    """Single-column Lunar content slide."""
    slide = _blank(prs)
    _bg(slide, LUNAR_50)
    _accent_bar(slide)

    if eyebrow:
        _eyebrow(slide, eyebrow)

    if title:
        _box(slide, title,
             MARGIN_L, Inches(0.75), Inches(11.5), Inches(0.75),
             size=24, bold=True, colour=BLACKHOLE)

    if body:
        _box(slide, body,
             MARGIN_L, CONTENT_T, CONTENT_W, CONTENT_H,
             size=12, colour=BLACKHOLE)

    if note:
        _box(slide, note,
             MARGIN_L, Inches(6.8), CONTENT_W, Inches(0.4),
             size=7, colour=LUNAR_700)


def make_demo(prs, team: str = "", title: str = "What we shipped  /  What we're building"):
    """Two-column demo slide: Shipped (left) | Building (right)."""
    slide = _blank(prs)
    _bg(slide, LUNAR_50)
    _accent_bar(slide)

    if team:
        _eyebrow(slide, team)

    _box(slide, title,
         MARGIN_L, Inches(0.75), Inches(11.5), Inches(0.75),
         size=24, bold=True, colour=BLACKHOLE)

    col_w = Inches(5.8)
    col_t = Inches(1.7)

    _box(slide, "Shipped",
         MARGIN_L, col_t, col_w, Inches(0.45),
         size=12, bold=True, colour=BLACKHOLE)
    _box(slide, "•\n•\n•",
         MARGIN_L, Inches(2.2), col_w, Inches(4.0),
         size=12, colour=BLACKHOLE)

    _box(slide, "Building",
         Inches(6.9), col_t, col_w, Inches(0.45),
         size=12, bold=True, colour=BLACKHOLE)
    _box(slide, "•\n•\n•",
         Inches(6.9), Inches(2.2), col_w, Inches(4.0),
         size=12, colour=BLACKHOLE)

    _divider_line(slide, Inches(6.55), col_t, Inches(5.3))


def make_qa(prs, team: str = "", title: str = "Q&A",
            note: str = ""):
    """Q&A slide with open notes area."""
    slide = _blank(prs)
    _bg(slide, LUNAR_50)
    _accent_bar(slide)

    if team:
        _eyebrow(slide, team)

    _box(slide, title,
         MARGIN_L, Inches(0.75), Inches(11.5), Inches(0.75),
         size=24, bold=True, colour=BLACKHOLE)

    if note:
        _box(slide, note,
             MARGIN_L, Inches(1.6), Inches(10), Inches(0.5),
             size=12, colour=LUNAR_700)

    _box(slide, "[Notes]",
         MARGIN_L, Inches(2.4), CONTENT_W, Inches(4.5),
         size=12, colour=RGBColor(0xC0, 0xBE, 0xBA))


def make_two_col(prs, eyebrow: str = "", title: str = "",
                 left_head: str = "", left_body: str = "",
                 right_head: str = "", right_body: str = ""):
    """Generic two-column content slide."""
    slide = _blank(prs)
    _bg(slide, LUNAR_50)
    _accent_bar(slide)

    if eyebrow:
        _eyebrow(slide, eyebrow)
    if title:
        _box(slide, title,
             MARGIN_L, Inches(0.75), Inches(11.5), Inches(0.75),
             size=24, bold=True, colour=BLACKHOLE)

    col_w = Inches(5.8)
    col_t = Inches(1.7)

    if left_head:
        _box(slide, left_head, MARGIN_L, col_t, col_w, Inches(0.45),
             size=12, bold=True, colour=BLACKHOLE)
    if left_body:
        _box(slide, left_body, MARGIN_L, Inches(2.2), col_w, Inches(4.5),
             size=12, colour=BLACKHOLE)

    if right_head:
        _box(slide, right_head, Inches(6.9), col_t, col_w, Inches(0.45),
             size=12, bold=True, colour=BLACKHOLE)
    if right_body:
        _box(slide, right_body, Inches(6.9), Inches(2.2), col_w, Inches(4.5),
             size=12, colour=BLACKHOLE)

    _divider_line(slide, Inches(6.55), col_t, Inches(5.3))


# ── Save ─────────────────────────────────────────────────────────────────────

def save_and_report(prs: Presentation, path: str):
    prs.save(path)
    count = len(prs.slides)
    print(f"Saved: {path}  ({count} slides)")
    print()
    print("To open in Google Slides:")
    print("  1. Go to drive.google.com")
    print("  2. Drag the file from Finder → drop into Drive")
    print("  3. Right-click → Open with → Google Slides")
